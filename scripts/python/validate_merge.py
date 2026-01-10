#!/usr/bin/env python3
"""
Validate merged PPTX output against source files for consistency.

Usage:
  python validate_merge.py <merged_output.pptx> <source1.pptx> <source2.pptx> ...
  python validate_merge.py --config merge-config.json
"""
import argparse
import sys
import os
from pptx import Presentation

# Add parent directory to path to import theme_resolver
sys.path.insert(0, os.path.dirname(__file__))
from theme_resolver import get_theme_color_rgb


def validate_slide_count(merged_path, source_paths):
    """Verify total slide count matches expectation."""
    merged = Presentation(merged_path)
    expected_slides = sum(len(Presentation(src).slides) for src in source_paths)
    actual_slides = len(merged.slides)
    
    status = "✅" if actual_slides == expected_slides else "❌"
    print(f"\n{status} Slide Count:")
    print(f"   Expected: {expected_slides} slides")
    print(f"   Actual:   {actual_slides} slides")
    return actual_slides == expected_slides


def validate_shapes(merged_path, source_paths):
    """Verify shapes were copied correctly."""
    merged = Presentation(merged_path)
    all_sources = [Presentation(src) for src in source_paths]
    
    issues = []
    slide_idx = 0
    
    for src_prs in all_sources:
        for src_slide in src_prs.slides:
            if slide_idx >= len(merged.slides):
                issues.append(f"Slide {slide_idx}: Missing in merged output")
                continue
            
            merged_slide = merged.slides[slide_idx]
            src_shape_count = len(src_slide.shapes)
            merged_shape_count = len(merged_slide.shapes)
            
            if src_shape_count != merged_shape_count:
                issues.append(f"Slide {slide_idx}: Shape count mismatch (src: {src_shape_count}, merged: {merged_shape_count})")
            
            slide_idx += 1
    
    status = "✅" if not issues else "❌"
    print(f"\n{status} Shapes:")
    if issues:
        for issue in issues[:5]:  # Show first 5 issues
            print(f"   {issue}")
        if len(issues) > 5:
            print(f"   ... and {len(issues) - 5} more issues")
    else:
        print(f"   All {slide_idx} slides have correct shape counts")
    
    return len(issues) == 0


def validate_text_content(merged_path, source_paths):
    """Verify text content was preserved."""
    merged = Presentation(merged_path)
    all_sources = [Presentation(src) for src in source_paths]
    
    issues = []
    slide_idx = 0
    
    for src_prs in all_sources:
        for src_slide in src_prs.slides:
            if slide_idx >= len(merged.slides):
                continue
            
            merged_slide = merged.slides[slide_idx]
            
            for src_shape, merged_shape in zip(src_slide.shapes, merged_slide.shapes):
                if hasattr(src_shape, 'text') and hasattr(merged_shape, 'text'):
                    if src_shape.text.strip() != merged_shape.text.strip():
                        issues.append(f"Slide {slide_idx}, Shape {src_shape.name}: Text mismatch")
            
            slide_idx += 1
    
    status = "✅" if not issues else "❌"
    print(f"\n{status} Text Content:")
    if issues:
        for issue in issues[:5]:
            print(f"   {issue}")
        if len(issues) > 5:
            print(f"   ... and {len(issues) - 5} more issues")
    else:
        print(f"   All text content preserved correctly")
    
    return len(issues) == 0


def validate_images(merged_path, source_paths):
    """Verify images were copied with correct data."""
    merged = Presentation(merged_path)
    all_sources = [Presentation(src) for src in source_paths]
    
    issues = []
    total_images = 0
    slide_idx = 0
    
    for src_prs in all_sources:
        for src_slide in src_prs.slides:
            if slide_idx >= len(merged.slides):
                continue
            
            merged_slide = merged.slides[slide_idx]
            
            for src_shape in src_slide.shapes:
                if src_shape.shape_type == 13:  # PICTURE
                    total_images += 1
                    # Find corresponding shape in merged
                    merged_shape = None
                    for ms in merged_slide.shapes:
                        if ms.name == src_shape.name and ms.shape_type == 13:
                            merged_shape = ms
                            break
                    
                    if not merged_shape:
                        issues.append(f"Slide {slide_idx}: Image {src_shape.name} missing")
                    else:
                        try:
                            src_size = len(src_shape.image.blob)
                            merged_size = len(merged_shape.image.blob)
                            if src_size != merged_size:
                                issues.append(f"Slide {slide_idx}: Image {src_shape.name} size mismatch")
                        except Exception as e:
                            issues.append(f"Slide {slide_idx}: Image {src_shape.name} error: {e}")
            
            slide_idx += 1
    
    status = "✅" if not issues else "❌"
    print(f"\n{status} Images:")
    print(f"   Total images: {total_images}")
    if issues:
        for issue in issues[:5]:
            print(f"   {issue}")
        if len(issues) > 5:
            print(f"   ... and {len(issues) - 5} more issues")
    else:
        print(f"   All images copied correctly")
    
    return len(issues) == 0


def validate_formatting(merged_path, source_paths):
    """Verify formatting (sizes, colors) was preserved."""
    merged = Presentation(merged_path)
    all_sources = [Presentation(src) for src in source_paths]
    
    issues = []
    total_runs = 0
    explicit_sizes = 0
    scheme_colors = 0
    rgb_colors = 0
    slide_idx = 0
    
    for src_prs in all_sources:
        for src_slide in src_prs.slides:
            if slide_idx >= len(merged.slides):
                continue
            
            merged_slide = merged.slides[slide_idx]
            
            for merged_shape in merged_slide.shapes:
                if hasattr(merged_shape, 'has_text_frame') and merged_shape.has_text_frame:
                    for para in merged_shape.text_frame.paragraphs:
                        for run in para.runs:
                            total_runs += 1
                            
                            # Check font size
                            if run.font.size:
                                explicit_sizes += 1
                            
                            # Check color type
                            try:
                                if run.font.color.type == 1:  # RGB
                                    rgb_colors += 1
                                elif run.font.color.type == 2:  # SCHEME
                                    scheme_colors += 1
                                    issues.append(f"Slide {slide_idx}: SCHEME color not converted to RGB")
                            except:
                                pass
            
            slide_idx += 1
    
    status = "✅" if scheme_colors == 0 else "⚠️"
    print(f"\n{status} Formatting:")
    print(f"   Total text runs: {total_runs}")
    print(f"   Explicit sizes: {explicit_sizes}/{total_runs} ({100*explicit_sizes/total_runs if total_runs > 0 else 0:.1f}%)")
    print(f"   RGB colors: {rgb_colors}")
    print(f"   SCHEME colors: {scheme_colors} {'(should be 0)' if scheme_colors > 0 else ''}")
    
    return scheme_colors == 0


def main():
    parser = argparse.ArgumentParser(description="Validate merged PPTX output")
    parser.add_argument('merged', help="Path to merged output PPTX file")
    parser.add_argument('sources', nargs='*', help="Source PPTX files that were merged")
    parser.add_argument('--config', '-c', help="Config file used for merge (to determine sources)")
    args = parser.parse_args()
    
    if not os.path.exists(args.merged):
        print(f"Error: Merged file not found: {args.merged}")
        sys.exit(1)
    
    source_paths = args.sources
    
    # If config provided, parse it to get sources
    if args.config:
        import json
        config_path = args.config
        if not os.path.isabs(config_path):
            # Try relative to assets/config
            candidate = os.path.join(os.path.dirname(__file__), "..", "..", "assets", "config", config_path)
            if os.path.exists(candidate):
                config_path = candidate
        
        with open(config_path) as f:
            config = json.load(f)
        
        # Extract unique source files
        source_set = set()
        for part in config.get('parts', []):
            if isinstance(part, str):
                source_set.add(part)
            else:
                source_set.add(part['file'])
        source_paths = list(source_set)
    
    if not source_paths:
        print("Error: No source files specified. Provide source paths or use --config")
        sys.exit(1)
    
    print("=" * 70)
    print("MERGE VALIDATION REPORT")
    print("=" * 70)
    print(f"Merged file: {args.merged}")
    print(f"Source files: {len(source_paths)}")
    for src in source_paths:
        print(f"  - {src}")
    print("=" * 70)
    
    # Run all validations
    results = []
    results.append(("Slide Count", validate_slide_count(args.merged, source_paths)))
    results.append(("Shapes", validate_shapes(args.merged, source_paths)))
    results.append(("Text Content", validate_text_content(args.merged, source_paths)))
    results.append(("Images", validate_images(args.merged, source_paths)))
    results.append(("Formatting", validate_formatting(args.merged, source_paths)))
    
    # Summary
    print("\n" + "=" * 70)
    print("SUMMARY:")
    print("=" * 70)
    for name, passed in results:
        status = "✅ PASS" if passed else "❌ FAIL"
        print(f"  {status}: {name}")
    
    all_passed = all(r[1] for r in results)
    print("=" * 70)
    if all_passed:
        print("✅ ALL VALIDATIONS PASSED")
        sys.exit(0)
    else:
        print("❌ SOME VALIDATIONS FAILED")
        sys.exit(1)


if __name__ == "__main__":
    main()
