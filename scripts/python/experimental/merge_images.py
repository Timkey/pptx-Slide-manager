#!/usr/bin/env python3
"""Merge PPTX files by converting every slide to an image and composing
those images into a new PPTX. This guarantees visual fidelity but makes
slides non-editable.

Requires `soffice` (LibreOffice) installed for conversion. The script
will detect `soffice` or `libreoffice` on PATH and use it for slide-to-PNG
conversion.

Usage:
  python3 scripts/python/merge_images.py [output.pptx] input1.pptx input2.pptx ...
  Output filename is optional - auto-generates if omitted.
"""
import argparse
import os
import random
import shlex
import shutil
import subprocess
import sys
import tempfile
from pptx import Presentation

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
OUTPUT_DIR = os.path.join(ROOT, "assets", "output", "merge_images")


def generate_output_filename(prefix="merged-images", extension=".pptx"):
    """Generate incremental filename in assets/output."""
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    counter = 1
    while True:
        filename = f"{prefix}-{counter:03d}{extension}"
        filepath = os.path.join(OUTPUT_DIR, filename)
        if not os.path.exists(filepath):
            return filepath
        counter += 1
        if counter > 999:
            filename = f"{prefix}-{random.randint(1000, 9999)}{extension}"
            return os.path.join(OUTPUT_DIR, filename)


def find_soffice():
    # Check common names on PATH first
    for cmd in ("soffice", "libreoffice"):
        path = shutil.which(cmd)
        if path:
            return path

    # On macOS, LibreOffice may be installed as an app without adding soffice to PATH.
    # Check the standard application bundle location.
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_path):
        return mac_path

    return None


def convert_to_png(soffice_cmd, src, outdir):
    # Use LibreOffice headless convert-to png
    cmd = [soffice_cmd, "--headless", "--convert-to", "png", "--outdir", outdir, src]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


def build_presentation_from_images(image_paths, output_path):
    prs = Presentation()
    # try to use a blank layout if available
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for img in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Fit the image to the slide size
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)


def main():
    p = argparse.ArgumentParser(description="Merge PPTX files by rendering slides to images")
    p.add_argument("-o", "--output", help="Output filename (saved to assets/output). If omitted, auto-generates name.")
    p.add_argument("inputs", nargs="+", help="Input PPTX files to convert and merge")
    args = p.parse_args()

    soffice = find_soffice()
    if not soffice:
        print("Error: LibreOffice (`soffice` or `libreoffice`) not found on PATH. Install LibreOffice to use this option.")
        sys.exit(1)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # Generate output filename if not provided
    if args.output:
        output_path = os.path.join(OUTPUT_DIR, args.output)
    else:
        output_path = generate_output_filename("merged-images", ".pptx")

    image_list = []
    with tempfile.TemporaryDirectory() as tmp:
        for src in args.inputs:
            if not os.path.isfile(src):
                print(f"Warning: input not found: {src}")
                continue
            convert_to_png(soffice, src, tmp)
            # LibreOffice emits files named after the source base; collect them
            base = os.path.splitext(os.path.basename(src))[0]
            # gather files that start with base and end with .png, sort to keep slide order
            matches = [os.path.join(tmp, f) for f in os.listdir(tmp) if f.startswith(base) and f.lower().endswith('.png')]
            matches.sort()
            image_list.extend(matches)

        if not image_list:
            print("No images produced; aborting.")
            sys.exit(1)

        build_presentation_from_images(image_list, output_path)
        print(f"Saved merged image-based presentation to: {output_path}")


if __name__ == "__main__":
    main()
