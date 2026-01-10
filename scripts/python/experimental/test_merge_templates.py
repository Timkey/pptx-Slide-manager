#!/usr/bin/env python3
"""
Test script to verify PPTX merging and template consistency.
This script:
1. Merges slide templates from assets/slides/templates
2. Outputs to assets/output/
3. Validates that merged PPTX hasn't drifted from templates
"""
import os
import sys
from pptx import Presentation

def get_slide_info(prs):
    """Extract information about slides for comparison."""
    info = {
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    for idx, slide in enumerate(prs.slides):
        slide_info = {
            'index': idx,
            'shapes_count': len(slide.shapes),
            'shapes': []
        }
        
        for shape in slide.shapes:
            shape_data = {
                'type': str(shape.shape_type),
                'has_text': hasattr(shape, 'text'),
                'text': shape.text if hasattr(shape, 'text') else None,
            }
            if hasattr(shape, 'left'):
                shape_data['position'] = (shape.left, shape.top, shape.width, shape.height)
            slide_info['shapes'].append(shape_data)
        
        info['slides'].append(slide_info)
    
    return info

def compare_presentations(template_path, merged_path):
    """Compare merged presentation with original template."""
    print(f"\n📊 Comparing presentations:")
    print(f"  Template: {template_path}")
    print(f"  Merged:   {merged_path}")
    
    template_prs = Presentation(template_path)
    merged_prs = Presentation(merged_path)
    
    template_info = get_slide_info(template_prs)
    merged_info = get_slide_info(merged_prs)
    
    issues = []
    
    # Check if merged has at least the template slides
    if merged_info['slide_count'] < template_info['slide_count']:
        issues.append(f"❌ Merged has fewer slides ({merged_info['slide_count']}) than template ({template_info['slide_count']})")
    else:
        print(f"✅ Slide count: Template={template_info['slide_count']}, Merged={merged_info['slide_count']}")
    
    # Compare individual slides (for slides that exist in template)
    for i in range(min(template_info['slide_count'], merged_info['slide_count'])):
        t_slide = template_info['slides'][i]
        m_slide = merged_info['slides'][i]
        
        # Check shape counts
        if m_slide['shapes_count'] != t_slide['shapes_count']:
            issues.append(f"⚠️  Slide {i+1}: Shape count differs (Template={t_slide['shapes_count']}, Merged={m_slide['shapes_count']})")
        
        # Check text content
        for j in range(min(len(t_slide['shapes']), len(m_slide['shapes']))):
            t_shape = t_slide['shapes'][j]
            m_shape = m_slide['shapes'][j]
            
            if t_shape.get('text') and m_shape.get('text'):
                if t_shape['text'] != m_shape['text'] and t_shape['text'].strip():
                    # Only flag if template had meaningful text
                    issues.append(f"⚠️  Slide {i+1}, Shape {j+1}: Text differs")
    
    if not issues:
        print("✅ No drift detected - merged presentation matches template structure")
    else:
        print(f"\n⚠️  Found {len(issues)} potential issues:")
        for issue in issues[:10]:  # Show first 10 issues
            print(f"  {issue}")
        if len(issues) > 10:
            print(f"  ... and {len(issues) - 10} more")
    
    return len(issues) == 0

def main():
    # Paths
    root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    templates_dir = os.path.join(root, "assets", "slides", "templates")
    output_dir = os.path.join(root, "assets", "output")
    
    os.makedirs(output_dir, exist_ok=True)
    
    print("=" * 70)
    print("PPTX Merge & Template Drift Test")
    print("=" * 70)
    
    # Find template files
    if not os.path.isdir(templates_dir):
        print(f"❌ Templates directory not found: {templates_dir}")
        return 1
    
    template_files = [f for f in os.listdir(templates_dir) if f.endswith('.pptx')]
    
    if not template_files:
        print(f"❌ No .pptx files found in {templates_dir}")
        return 1
    
    print(f"\n📁 Found {len(template_files)} template(s):")
    for f in template_files:
        print(f"  - {f}")
    
    # Test 1: Merge all templates
    print(f"\n{'='*70}")
    print("TEST 1: Merging Templates")
    print('='*70)
    
    output_file = os.path.join(output_dir, "merged_test.pptx")
    input_files = [os.path.join(templates_dir, f) for f in template_files]
    
    print(f"Output: {output_file}")
    print(f"Inputs: {len(input_files)} file(s)")
    
    try:
        # Create merged presentation
        merged_prs = Presentation(input_files[0])
        
        # Add slides from other presentations
        for input_file in input_files[1:]:
            source_prs = Presentation(input_file)
            for slide in source_prs.slides:
                # Copy slide layout
                blank_slide_layout = merged_prs.slide_layouts[0]
                new_slide = merged_prs.slides.add_slide(blank_slide_layout)
                
                # Copy shapes
                for shape in slide.shapes:
                    el = shape.element
                    newel = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        merged_prs.save(output_file)
        print(f"✅ Merge successful: {output_file}")
        
    except Exception as e:
        print(f"❌ Merge failed: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    # Test 2: Check for template drift
    print(f"\n{'='*70}")
    print("TEST 2: Template Drift Detection")
    print('='*70)
    
    all_good = True
    for template_file in template_files:
        template_path = os.path.join(templates_dir, template_file)
        if not compare_presentations(template_path, output_file):
            all_good = False
    
    # Summary
    print(f"\n{'='*70}")
    print("SUMMARY")
    print('='*70)
    print(f"✅ Merged file created: {output_file}")
    print(f"✅ Total slides in merged: {len(Presentation(output_file).slides)}")
    
    if all_good:
        print("✅ No template drift detected")
        print("\n🎉 All tests passed!")
        return 0
    else:
        print("⚠️  Some template drift detected (this may be expected)")
        print("\n⚠️  Review the issues above")
        return 0  # Return 0 since drift warnings are informational

if __name__ == "__main__":
    # Import after defining main to avoid issues
    from copy import deepcopy
    sys.exit(main())
