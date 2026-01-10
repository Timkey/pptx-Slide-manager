from pptx import Presentation
from copy import deepcopy

# Load source
src = Presentation('/app/assets/slides/templates/test1.pptx')
print(f"Source (test1.pptx) has {len(src.slides)} slide(s)")
for idx, slide in enumerate(src.slides, 1):
    print(f"  Slide {idx}: {len(slide.shapes)} shapes")
    for s in slide.shapes:
        print(f"    - {s.shape_type}: {getattr(s, 'name', 'unnamed')}")

# Try to copy
base = Presentation('/app/assets/slides/templates/test1.pptx')
# Remove slides
while len(base.slides) > 0:
    rId = base.slides._sldIdLst[0].rId
    base.part.drop_rel(rId)
    del base.slides._sldIdLst[0]

print(f"\nAfter clearing: {len(base.slides)} slides")

# Now add back
for slide in src.slides:
    layout = base.slide_layouts[0]
    target = base.slides.add_slide(layout)
    spTree = target.shapes._spTree
    for shape in slide.shapes:
        try:
            el = deepcopy(shape._element)
            spTree.append(el)
            print(f"  ✓ Copied {getattr(shape, 'name', 'unnamed')}")
        except Exception as e:
            print(f"  ✗ Failed to copy {getattr(shape, 'name', 'unnamed')}: {e}")

print(f"\nResult: {len(base.slides)} slides")
for idx, slide in enumerate(base.slides, 1):
    print(f"  Slide {idx}: {len(slide.shapes)} shapes")
