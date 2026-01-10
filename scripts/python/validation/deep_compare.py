from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

def analyze_text_properties(shape, shape_name):
    """Extract detailed text properties from a shape."""
    if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
        return None
    
    details = {
        'shape_name': shape_name,
        'text': shape.text,
        'paragraphs': []
    }
    
    for para_idx, para in enumerate(shape.text_frame.paragraphs):
        para_info = {
            'text': para.text,
            'level': para.level,
            'alignment': str(para.alignment) if para.alignment else 'None',
            'has_bullet': False,
            'runs': []
        }
        
        # Check bullet/numbering
        if para.font.size:
            para_info['bullet'] = 'has bullet' if para.level > 0 else 'no bullet'
        
        for run_idx, run in enumerate(para.runs):
            run_info = {
                'text': run.text,
                'font_name': run.font.name,
                'font_size': run.font.size.pt if run.font.size else None,
                'bold': run.font.bold,
                'italic': run.font.italic,
                'underline': run.font.underline,
            }
            
            # Color
            if run.font.color.type:
                if run.font.color.type == 1:  # RGB
                    run_info['color'] = f'RGB({run.font.color.rgb})'
                elif run.font.color.type == 2:  # SCHEME
                    run_info['color'] = f'SCHEME({run.font.color.theme_color})'
                else:
                    run_info['color'] = f'TYPE({run.font.color.type})'
            else:
                run_info['color'] = 'None'
            
            para_info['runs'].append(run_info)
        
        details['paragraphs'].append(para_info)
    
    return details

def analyze_shape_properties(shape):
    """Extract shape properties."""
    props = {
        'type': str(shape.shape_type),
        'name': shape.name,
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
    }
    
    # Fill color
    if hasattr(shape, 'fill'):
        if shape.fill.type:
            props['fill_type'] = str(shape.fill.type)
            if shape.fill.type == 1:  # SOLID
                if shape.fill.fore_color.type == 1:  # RGB
                    props['fill_color'] = f'RGB({shape.fill.fore_color.rgb})'
                elif shape.fill.fore_color.type == 2:  # SCHEME
                    props['fill_color'] = f'SCHEME({shape.fill.fore_color.theme_color})'
        else:
            props['fill_type'] = 'None'
    
    # Line color
    if hasattr(shape, 'line'):
        if shape.line.color.type:
            if shape.line.color.type == 1:
                props['line_color'] = f'RGB({shape.line.color.rgb})'
            elif shape.line.color.type == 2:
                props['line_color'] = f'SCHEME({shape.line.color.theme_color})'
    
    return props

def analyze_picture(shape):
    """Extract picture properties."""
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    
    try:
        return {
            'image_size': len(shape.image.blob),
            'content_type': shape.image.content_type,
            'accessible': True
        }
    except:
        return {'accessible': False}

print('='*80)
print('DEEP CONTENT COMPARISON')
print('='*80)

# Analyze source files
sources = {}
for fname in ['test1.pptx', 'test2.pptx']:
    prs = Presentation(f'/app/assets/slides/templates/{fname}')
    sources[fname] = []
    
    print(f'\n📁 SOURCE: {fname}')
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f'\n  Slide {slide_idx}:')
        slide_data = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_data = {
                'index': shape_idx,
                'properties': analyze_shape_properties(shape)
            }
            
            print(f'    Shape {shape_idx}: {shape.name} ({shape.shape_type})')
            
            # Position and size
            print(f'      Position: ({shape.left}, {shape.top}), Size: ({shape.width}, {shape.height})')
            
            # Text content
            text_props = analyze_text_properties(shape, shape.name)
            if text_props:
                shape_data['text'] = text_props
                print(f'      Text: "{shape.text}"')
                for para_idx, para in enumerate(text_props['paragraphs']):
                    print(f'        Para {para_idx}: level={para["level"]}, align={para["alignment"]}')
                    print(f'          Text: "{para["text"]}"')
                    for run_idx, run in enumerate(para['runs']):
                        print(f'          Run {run_idx}: "{run["text"]}"')
                        print(f'            Font: {run["font_name"]}, Size: {run["font_size"]}, Color: {run["color"]}')
                        print(f'            Bold: {run["bold"]}, Italic: {run["italic"]}')
            
            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_props = analyze_picture(shape)
                shape_data['picture'] = pic_props
                if pic_props and pic_props.get('accessible'):
                    print(f'      Picture: {pic_props["image_size"]} bytes, {pic_props["content_type"]}')
                else:
                    print(f'      Picture: ❌ NOT ACCESSIBLE')
            
            slide_data.append(shape_data)
        
        sources[fname].append(slide_data)

# Analyze merged file
print('\n' + '='*80)
print('📄 MERGED OUTPUT: merged-pptx-001.pptx')
print('='*80)

prs = Presentation('/app/assets/output/merge_pptx/merged-pptx-001.pptx')
merged_data = []

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f'\n  Slide {slide_idx}:')
    slide_data = []
    
    for shape_idx, shape in enumerate(slide.shapes):
        shape_data = {
            'index': shape_idx,
            'properties': analyze_shape_properties(shape)
        }
        
        print(f'    Shape {shape_idx}: {shape.name} ({shape.shape_type})')
        print(f'      Position: ({shape.left}, {shape.top}), Size: ({shape.width}, {shape.height})')
        
        text_props = analyze_text_properties(shape, shape.name)
        if text_props:
            shape_data['text'] = text_props
            print(f'      Text: "{shape.text}"')
            for para_idx, para in enumerate(text_props['paragraphs']):
                print(f'        Para {para_idx}: level={para["level"]}, align={para["alignment"]}')
                print(f'          Text: "{para["text"]}"')
                for run_idx, run in enumerate(para['runs']):
                    print(f'          Run {run_idx}: "{run["text"]}"')
                    print(f'            Font: {run["font_name"]}, Size: {run["font_size"]}, Color: {run["color"]}')
                    print(f'            Bold: {run["bold"]}, Italic: {run["italic"]}')
        
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_props = analyze_picture(shape)
            shape_data['picture'] = pic_props
            if pic_props and pic_props.get('accessible'):
                print(f'      Picture: {pic_props["image_size"]} bytes, {pic_props["content_type"]}')
            else:
                print(f'      Picture: ❌ NOT ACCESSIBLE')
        
        slide_data.append(shape_data)
    
    merged_data.append(slide_data)

# Compare
print('\n' + '='*80)
print('🔍 DIFFERENCES DETECTED')
print('='*80)

differences = []

# Compare slide 1 (test1)
print('\n📊 Slide 1 Comparison (from test1.pptx):')
source_slide = sources['test1.pptx'][0]
merged_slide = merged_data[0]

for idx in range(max(len(source_slide), len(merged_slide))):
    if idx >= len(source_slide):
        print(f'  ❌ Extra shape {idx} in merged (not in source)')
        differences.append(f'Slide 1: Extra shape {idx}')
        continue
    if idx >= len(merged_slide):
        print(f'  ❌ Missing shape {idx} in merged (was in source)')
        differences.append(f'Slide 1: Missing shape {idx}')
        continue
    
    src = source_slide[idx]
    mrg = merged_slide[idx]
    
    # Compare text
    if 'text' in src and 'text' in mrg:
        if src['text']['text'] != mrg['text']['text']:
            print(f'  ❌ Shape {idx}: Text mismatch')
            print(f'      Source: "{src["text"]["text"]}"')
            print(f'      Merged: "{mrg["text"]["text"]}"')
            differences.append(f'Slide 1 Shape {idx}: Text content differs')
        
        # Compare paragraph properties
        for para_idx in range(max(len(src['text']['paragraphs']), len(mrg['text']['paragraphs']))):
            if para_idx >= len(src['text']['paragraphs']) or para_idx >= len(mrg['text']['paragraphs']):
                continue
            
            src_para = src['text']['paragraphs'][para_idx]
            mrg_para = mrg['text']['paragraphs'][para_idx]
            
            if src_para['level'] != mrg_para['level']:
                print(f'  ❌ Shape {idx} Para {para_idx}: Level mismatch (src={src_para["level"]}, merged={mrg_para["level"]})')
                differences.append(f'Slide 1 Shape {idx} Para {para_idx}: Bullet level differs')
            
            if src_para['alignment'] != mrg_para['alignment']:
                print(f'  ❌ Shape {idx} Para {para_idx}: Alignment mismatch')
                differences.append(f'Slide 1 Shape {idx} Para {para_idx}: Alignment differs')
            
            # Compare runs
            for run_idx in range(max(len(src_para['runs']), len(mrg_para['runs']))):
                if run_idx >= len(src_para['runs']) or run_idx >= len(mrg_para['runs']):
                    continue
                
                src_run = src_para['runs'][run_idx]
                mrg_run = mrg_para['runs'][run_idx]
                
                if src_run['color'] != mrg_run['color']:
                    print(f'  ❌ Shape {idx} Para {para_idx} Run {run_idx}: Color mismatch')
                    print(f'      Source: {src_run["color"]}')
                    print(f'      Merged: {mrg_run["color"]}')
                    differences.append(f'Slide 1 Shape {idx}: Color differs')
                
                if src_run['font_size'] != mrg_run['font_size']:
                    print(f'  ❌ Shape {idx} Para {para_idx} Run {run_idx}: Font size mismatch ({src_run["font_size"]} vs {mrg_run["font_size"]})')
                    differences.append(f'Slide 1 Shape {idx}: Font size differs')

print('\n' + '='*80)
if differences:
    print(f'❌ FOUND {len(differences)} DIFFERENCES')
    for diff in differences:
        print(f'  • {diff}')
else:
    print('✅ NO DIFFERENCES FOUND - PERFECT MATCH')
print('='*80)
