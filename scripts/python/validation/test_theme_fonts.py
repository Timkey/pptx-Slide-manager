from pptx import Presentation

def get_theme_font_name(source_slide, is_major=False):
    """Get the theme font name (major for headings, minor for body) from source slide's master."""
    try:
        master = source_slide.slide_layout.slide_master
        master_elem = master._element
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        
        # Find themeElements
        themeElements = master_elem.find(f'.//{ns_a}themeElements')
        if themeElements is not None:
            fontScheme = themeElements.find(f'{ns_a}fontScheme')
            if fontScheme is not None:
                if is_major:
                    font = fontScheme.find(f'{ns_a}majorFont')
                else:
                    font = fontScheme.find(f'{ns_a}minorFont')
                
                if font is not None:
                    latin = font.find(f'{ns_a}latin')
                    if latin is not None:
                        return latin.get('typeface')
    except:
        pass
    return None

prs2 = Presentation('/app/assets/slides/templates/test2.pptx')
slide = prs2.slides[0]

major_font = get_theme_font_name(slide, is_major=True)
minor_font = get_theme_font_name(slide, is_major=False)

print(f'test2.pptx theme fonts:')
print(f'  Major (headings): {major_font}')
print(f'  Minor (body): {minor_font}')

prs1 = Presentation('/app/assets/slides/templates/test1.pptx')
slide = prs1.slides[0]

major_font = get_theme_font_name(slide, is_major=True)
minor_font = get_theme_font_name(slide, is_major=False)

print(f'\ntest1.pptx theme fonts:')
print(f'  Major (headings): {major_font}')
print(f'  Minor (body): {minor_font}')
