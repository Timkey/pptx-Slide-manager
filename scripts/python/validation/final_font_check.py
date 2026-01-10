from pptx import Presentation

src2 = Presentation('/app/assets/slides/templates/test2.pptx')
merged = Presentation('/app/assets/output/merge_pptx/merged-pptx-001.pptx')

print('FONT AND COLOR COMPARISON')
print('='*70)

# Get "Priest:" runs
src_slide = src2.slides[0]
merged_slide = merged.slides[1]

for src_shape, merged_shape in zip(src_slide.shapes, merged_slide.shapes):
    if hasattr(src_shape, 'has_text_frame') and src_shape.has_text_frame:
        for para_idx, (src_para, merged_para) in enumerate(zip(src_shape.text_frame.paragraphs, merged_shape.text_frame.paragraphs)):
            for run_idx, (src_run, merged_run) in enumerate(zip(src_para.runs, merged_para.runs)):
                if 'Priest' in src_run.text or 'People' in src_run.text:
                    print(f'\nPara {para_idx} Run {run_idx}: "{src_run.text}"')
                    
                    # Font name
                    src_font_name = src_run.font.name
                    merged_font_name = merged_run.font.name
                    print(f'  Font name: src={src_font_name}, merged={merged_font_name}')
                    
                    # Font size
                    src_size = src_run.font.size.pt if src_run.font.size else None
                    merged_size = merged_run.font.size.pt if merged_run.font.size else None
                    print(f'  Font size: src={src_size}, merged={merged_size}')
                    
                    # Color
                    try:
                        if src_run.font.color.type:
                            if src_run.font.color.type == 1:
                                print(f'  Source color: RGB({src_run.font.color.rgb})')
                            elif src_run.font.color.type == 2:
                                print(f'  Source color: SCHEME({src_run.font.color.theme_color})')
                    except:
                        print(f'  Source color: None')
                    
                    try:
                        if merged_run.font.color.type:
                            if merged_run.font.color.type == 1:
                                print(f'  Merged color: RGB({merged_run.font.color.rgb})')
                            elif merged_run.font.color.type == 2:
                                print(f'  Merged color: SCHEME({merged_run.font.color.theme_color})')
                    except:
                        print(f'  Merged color: None')
                    
                    # Match check
                    fonts_match = src_font_name == merged_font_name
                    sizes_match = src_size == merged_size
                    
                    if fonts_match and sizes_match:
                        print(f'  ✅ Matches')
                    else:
                        print(f'  ❌ Mismatch!')

print('\n' + '='*70)
print('Please check the opened PowerPoint file to verify visually')
print('='*70)
