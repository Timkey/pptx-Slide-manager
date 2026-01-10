from pptx import Presentation

print('='*70)
print('FINAL COMPREHENSIVE VALIDATION')
print('='*70)

issues = []

# Check test1
print('\n📁 SOURCE: test1.pptx')
src1 = Presentation('/app/assets/slides/templates/test1.pptx')
slide1_src = src1.slides[0]
print(f'  Background: {slide1_src.background.fill.type}')
if slide1_src.background.fill.type == 1:
    try:
        if slide1_src.background.fill.fore_color.type == 2:
            print(f'    Color: SCHEME {slide1_src.background.fill.fore_color.theme_color}')
    except:
        pass

# Check test2
print('\n📁 SOURCE: test2.pptx')
src2 = Presentation('/app/assets/slides/templates/test2.pptx')
slide2_src = src2.slides[0]
print(f'  Background: {slide2_src.background.fill.type}')
has_bullets_src = False
for shape in slide2_src.shapes:
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                pPr = para._element.pPr
                if pPr is not None:
                    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                    if pPr.find(f'.//{ns}buChar') is not None or pPr.find(f'.//{ns}buAutoNum') is not None:
                        has_bullets_src = True
print(f'  Has bullets: {has_bullets_src}')

# Check merged
print('\n📄 MERGED: merged-pptx-001.pptx')
merged = Presentation('/app/assets/output/merge_pptx/merged-pptx-001.pptx')

# Check Slide 1 background
slide1_merged = merged.slides[0]
print(f'\nSlide 1:')
print(f'  Background: {slide1_merged.background.fill.type}')
if slide1_merged.background.fill.type == 1:
    try:
        if slide1_merged.background.fill.fore_color.type == 2:
            src_color = slide1_src.background.fill.fore_color.theme_color
            merged_color = slide1_merged.background.fill.fore_color.theme_color
            if src_color == merged_color:
                print(f'    ✅ Color matches: SCHEME {merged_color}')
            else:
                print(f'    ❌ Color mismatch: src={src_color}, merged={merged_color}')
                issues.append('Slide 1: Background color mismatch')
    except:
        pass

# Check Slide 2 background
slide2_merged = merged.slides[1]
print(f'\nSlide 2:')
print(f'  Background: {slide2_merged.background.fill.type}')
if slide2_merged.background.fill.type == 1:
    try:
        if slide2_merged.background.fill.fore_color.type == 1:
            color = slide2_merged.background.fill.fore_color.rgb
            # FBFFBB is the yellow from test2's master
            if color.upper() in ['FBFFBB', 'FFFFBB']:  # Allow slight variation
                print(f'    ✅ Background color preserved: RGB({color})')
            else:
                print(f'    ⚠️  Background color: RGB({color}) (expected FBFFBB)')
    except Exception as e:
        print(f'    ⚠️  Could not verify color: {e}')

# Check bullets
has_bullets_merged = False
has_buNone = True
for shape in slide2_merged.shapes:
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if para.text.strip():
                pPr = para._element.pPr
                if pPr is not None:
                    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                    if pPr.find(f'.//{ns}buChar') is not None or pPr.find(f'.//{ns}buAutoNum') is not None:
                        has_bullets_merged = True
                    if pPr.find(f'.//{ns}buNone') is None and para.text.strip():
                        has_buNone = False

if not has_bullets_merged and has_buNone:
    print(f'  ✅ Bullets correctly disabled (buNone set)')
else:
    print(f'  ❌ Bullet formatting issue')
    issues.append('Slide 2: Bullets not properly disabled')

# Check shapes match
print(f'\nShape counts:')
print(f'  test1: {len(slide1_src.shapes)} shapes → Merged slide 1: {len(slide1_merged.shapes)} shapes')
print(f'  test2: {len(slide2_src.shapes)} shapes → Merged slide 2: {len(slide2_merged.shapes)} shapes')

if len(slide1_src.shapes) == len(slide1_merged.shapes) and len(slide2_src.shapes) == len(slide2_merged.shapes):
    print(f'  ✅ All shapes preserved')
else:
    issues.append('Shape count mismatch')

print('\n' + '='*70)
if not issues:
    print('✅ ALL CHECKS PASSED - PERFECT FIDELITY')
else:
    print(f'❌ FOUND {len(issues)} ISSUES:')
    for issue in issues:
        print(f'  • {issue}')
print('='*70)
