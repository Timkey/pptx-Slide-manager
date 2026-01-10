from pptx import Presentation

print('='*70)
print('FONT SIZE VALIDATION')
print('='*70)

issues = []

# Check test2 source
src2 = Presentation('/app/assets/slides/templates/test2.pptx')
merged = Presentation('/app/assets/output/merge_pptx/merged-pptx-001.pptx')

print('\n📊 Comparing test2.pptx Slide 1 with merged Slide 2:')
print('Expected: Body text should be 28pt (from test2 master)')

src_slide = src2.slides[0]
merged_slide = merged.slides[1]

for shape_idx, (src_shape, merged_shape) in enumerate(zip(src_slide.shapes, merged_slide.shapes)):
    if hasattr(src_shape, 'has_text_frame') and src_shape.has_text_frame:
        print(f'\nShape {shape_idx}: {src_shape.name}')
        
        # Get placeholder type to determine expected default
        expected_size = None
        if src_shape.is_placeholder:
            try:
                ph_type = src_shape.placeholder_format.type
                if ph_type == 1:  # TITLE
                    expected_size = 44  # Title default from master
                elif ph_type == 2:  # BODY
                    expected_size = 28  # Body default from master
            except:
                pass
        
        for para_idx, (src_para, merged_para) in enumerate(zip(src_shape.text_frame.paragraphs, merged_shape.text_frame.paragraphs)):
            if src_para.text.strip():
                for run_idx, (src_run, merged_run) in enumerate(zip(src_para.runs, merged_para.runs)):
                    if src_run.text.strip():
                        src_size = src_run.font.size.pt if src_run.font.size else None
                        merged_size = merged_run.font.size.pt if merged_run.font.size else None
                        
                        # If source has no explicit size, it uses the master default
                        effective_src_size = src_size if src_size else expected_size
                        
                        # Check if merged matches
                        if merged_size:
                            if effective_src_size and abs(merged_size - effective_src_size) > 0.1:
                                print(f'  ❌ Para {para_idx} Run {run_idx}: "{src_run.text[:20]}"')
                                print(f'     Expected: {effective_src_size}pt, Got: {merged_size}pt')
                                issues.append(f'Font size mismatch: {merged_size} vs {effective_src_size}')
                            else:
                                print(f'  ✅ Para {para_idx} Run {run_idx}: {merged_size}pt')
                        else:
                            if effective_src_size:
                                print(f'  ❌ Para {para_idx} Run {run_idx}: No explicit size (will use wrong master)')
                                issues.append(f'Missing explicit font size')
                            else:
                                print(f'  ⚠️  Para {para_idx} Run {run_idx}: Both using defaults')

print('\n' + '='*70)
if not issues:
    print('✅ ALL FONT SIZES PRESERVED CORRECTLY')
else:
    print(f'❌ FOUND {len(issues)} ISSUES')
    for issue in set(issues[:5]):  # Show unique issues
        print(f'  • {issue}')
print('='*70)
