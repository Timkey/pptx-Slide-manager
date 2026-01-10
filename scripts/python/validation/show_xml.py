from pptx import Presentation

src2 = Presentation('/app/assets/slides/templates/test2.pptx')
merged = Presentation('/app/assets/output/merge_pptx/merged-pptx-001.pptx')

print('DETAILED XML COMPARISON - "Priest:" text')
print('='*70)

# Get the first "Priest:" run from test2
src_slide = src2.slides[0]
merged_slide = merged.slides[1]

for src_shape, merged_shape in zip(src_slide.shapes, merged_slide.shapes):
    if hasattr(src_shape, 'has_text_frame') and src_shape.has_text_frame:
        for src_para, merged_para in zip(src_shape.text_frame.paragraphs, merged_shape.text_frame.paragraphs):
            for src_run, merged_run in zip(src_para.runs, merged_para.runs):
                if 'Priest:' in src_run.text:
                    print(f'\nText: "{src_run.text}"')
                    print('\nSOURCE XML:')
                    from lxml import etree
                    print(etree.tostring(src_run.element, pretty_print=True, encoding='unicode')[:1000])
                    
                    print('\nMERGED XML:')
                    print(etree.tostring(merged_run.element, pretty_print=True, encoding='unicode')[:1000])
                    break
