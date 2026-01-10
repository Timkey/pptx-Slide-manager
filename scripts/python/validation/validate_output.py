#!/usr/bin/env python3
"""
Validate merged PPTX and PDF outputs for corruption and consistency.
Runs inside Docker container to ensure consistent environment.
"""
import os
import sys
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.exc import PackageNotFoundError

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
OUTPUT_DIR = os.path.join(ROOT, "assets", "output")


def analyze_slide_structure(slide):
    """Analyze the structure of a single slide."""
    structure = {
        "shapes": len(slide.shapes),
        "shape_types": {},
        "has_title": False,
        "has_text": False,
        "text_boxes": 0,
        "pictures": 0,
        "tables": 0,
        "charts": 0,
        "placeholders": 0,
        "groups": 0,
        "picture_sizes": [],  # Track image blob sizes
        "fonts": set(),       # Track fonts used
        "has_broken_images": False
    }
    
    for shape in slide.shapes:
        # Count shape types
        shape_type = str(shape.shape_type)
        structure["shape_types"][shape_type] = structure["shape_types"].get(shape_type, 0) + 1
        
        # Check for title
        if shape.has_text_frame:
            structure["has_text"] = True
            if hasattr(shape, "name") and "title" in shape.name.lower():
                structure["has_title"] = True
            
            # Collect fonts
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            structure["fonts"].add(run.font.name)
            except:
                pass
        
        # Count specific shape types
        if hasattr(shape, "shape_type"):
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                structure["text_boxes"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                structure["pictures"] += 1
                # Try to access the image to check if relationship is intact
                try:
                    blob = shape.image.blob
                    structure["picture_sizes"].append(len(blob))
                except Exception as e:
                    structure["has_broken_images"] = True
                    structure["picture_sizes"].append(0)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                structure["tables"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                structure["charts"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                structure["placeholders"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                structure["groups"] += 1
    
    return structure


def compare_slide_structures(source_structures, merged_structures):
    """Compare source and merged slide structures to detect issues.
    
    Note: This compares the total content - merged should have ALL slides
    from ALL source files in order.
    """
    issues = []
    
    if len(merged_structures) < len(source_structures):
        issues.append(f"Missing slides: expected {len(source_structures)}, found {len(merged_structures)}")
    elif len(merged_structures) > len(source_structures):
        issues.append(f"Extra slides: expected {len(source_structures)}, found {len(merged_structures)}")
    
    # Check total shape count across all slides
    total_src_shapes = sum(s["shapes"] for s in source_structures)
    total_merged_shapes = sum(s["shapes"] for s in merged_structures)
    
    if total_merged_shapes < total_src_shapes:
        issues.append(f"Lost shapes overall: {total_src_shapes} in sources → {total_merged_shapes} in merged")
    
    # Check for significant content loss on individual slides
    for idx, (src, merged) in enumerate(zip(source_structures, merged_structures), 1):
        # Only flag if we lost MORE than 1 shape (some minor differences are expected)
        if merged["shapes"] < src["shapes"] - 1:
            issues.append(f"Slide {idx}: Significant shape loss ({src['shapes']} → {merged['shapes']})")
        
        # Check for broken images
        if merged.get("has_broken_images"):
            issues.append(f"Slide {idx}: Broken image relationships (cannot access image data)")
        
        # Check for lost important content types
        if src["pictures"] > 0 and merged["pictures"] == 0:
            issues.append(f"Slide {idx}: All images lost ({src['pictures']} pictures missing)")
        elif src["pictures"] > 0 and merged["pictures"] > 0:
            # Check if image sizes match (indicates images were properly copied)
            src_sizes = sorted(src.get("picture_sizes", []))
            merged_sizes = sorted(merged.get("picture_sizes", []))
            if src_sizes and merged_sizes and src_sizes != merged_sizes:
                issues.append(f"Slide {idx}: Image data mismatch (sizes: {src_sizes} → {merged_sizes})")
        
        # Check for font loss
        src_fonts = src.get("fonts", set())
        merged_fonts = merged.get("fonts", set())
        if src_fonts and not merged_fonts:
            issues.append(f"Slide {idx}: All font information lost")
        elif src_fonts and merged_fonts:
            missing_fonts = src_fonts - merged_fonts
            if missing_fonts:
                issues.append(f"Slide {idx}: Font changes detected (lost: {', '.join(missing_fonts)})")
        
        if src["tables"] > 0 and merged["tables"] == 0:
            issues.append(f"Slide {idx}: All tables lost")
        
        if src["charts"] > 0 and merged["charts"] == 0:
            issues.append(f"Slide {idx}: All charts lost")
    
    return issues


def validate_pptx(filepath, source_templates=None):
    """Validate a PPTX file for corruption and consistency.
    
    Args:
        filepath: Path to PPTX file to validate
        source_templates: Optional list of source template paths for structure comparison
    """
    results = {
        "file": os.path.basename(filepath),
        "valid": False,
        "errors": [],
        "warnings": [],
        "info": {}
    }
    
    # Check 1: File exists and has size
    if not os.path.exists(filepath):
        results["errors"].append("File does not exist")
        return results
    
    file_size = os.path.getsize(filepath)
    if file_size == 0:
        results["errors"].append("File is empty (0 bytes)")
        return results
    
    results["info"]["size_bytes"] = file_size
    results["info"]["size_kb"] = round(file_size / 1024, 2)
    
    # Check 2: Valid ZIP structure (PPTX is a ZIP archive)
    try:
        with zipfile.ZipFile(filepath, 'r') as zf:
            # Check for required OOXML files
            namelist = zf.namelist()
            required_files = ['[Content_Types].xml', '_rels/.rels']
            missing = [f for f in required_files if f not in namelist]
            if missing:
                results["errors"].append(f"Missing required files: {missing}")
            
            # Check for corrupted ZIP entries
            bad_file = zf.testzip()
            if bad_file:
                results["errors"].append(f"Corrupted ZIP entry: {bad_file}")
            
            results["info"]["zip_entries"] = len(namelist)
            
    except zipfile.BadZipFile as e:
        results["errors"].append(f"Not a valid ZIP file: {e}")
        return results
    except Exception as e:
        results["errors"].append(f"ZIP validation error: {e}")
        return results
    
    # Check 3: Load with python-pptx
    try:
        prs = Presentation(filepath)
        results["info"]["slides"] = len(prs.slides)
        results["info"]["slide_layouts"] = len(prs.slide_layouts)
        results["info"]["slide_masters"] = len(prs.slide_masters)
        
        # Check for empty presentation
        if len(prs.slides) == 0:
            results["warnings"].append("Presentation has no slides")
        
        # Validate each slide can be accessed
        for idx, slide in enumerate(prs.slides, 1):
            try:
                _ = slide.shapes
                _ = slide.slide_layout
            except Exception as e:
                results["errors"].append(f"Slide {idx} is corrupted: {e}")
        
        # Check slide dimensions
        results["info"]["width"] = prs.slide_width
        results["info"]["height"] = prs.slide_height
        
        # Analyze slide structures
        slide_structures = []
        for idx, slide in enumerate(prs.slides, 1):
            try:
                structure = analyze_slide_structure(slide)
                slide_structures.append(structure)
                results["info"][f"slide_{idx}_shapes"] = structure["shapes"]
                
                # Check for common issues
                if structure["shapes"] == 0:
                    results["warnings"].append(f"Slide {idx} has no shapes (empty slide)")
                
                if not structure["has_text"] and structure["pictures"] == 0:
                    results["warnings"].append(f"Slide {idx} has no text or images")
            except Exception as e:
                results["warnings"].append(f"Slide {idx} structure analysis failed: {e}")
        
        # If source templates provided, compare structures
        if source_templates:
            source_structures = []
            for template_path in source_templates:
                if os.path.exists(template_path):
                    try:
                        src_prs = Presentation(template_path)
                        for slide in src_prs.slides:
                            source_structures.append(analyze_slide_structure(slide))
                    except Exception as e:
                        results["warnings"].append(f"Could not analyze source {os.path.basename(template_path)}: {e}")
            
            if source_structures:
                structure_issues = compare_slide_structures(source_structures, slide_structures)
                for issue in structure_issues:
                    results["errors"].append(f"Structure issue: {issue}")
        
    except PackageNotFoundError as e:
        results["errors"].append(f"PPTX package error: {e}")
        return results
    except Exception as e:
        results["errors"].append(f"PPTX loading error: {e}")
        return results
    
    # Mark as valid if no errors
    results["valid"] = len(results["errors"]) == 0
    
    return results


def validate_pdf(filepath):
    """Validate a PDF file for corruption."""
    results = {
        "file": os.path.basename(filepath),
        "valid": False,
        "errors": [],
        "warnings": [],
        "info": {}
    }
    
    # Check 1: File exists and has size
    if not os.path.exists(filepath):
        results["errors"].append("File does not exist")
        return results
    
    file_size = os.path.getsize(filepath)
    if file_size == 0:
        results["errors"].append("File is empty (0 bytes)")
        return results
    
    results["info"]["size_bytes"] = file_size
    results["info"]["size_kb"] = round(file_size / 1024, 2)
    
    # Check 2: PDF header
    try:
        with open(filepath, 'rb') as f:
            header = f.read(5)
            if not header.startswith(b'%PDF-'):
                results["errors"].append("Invalid PDF header")
            else:
                results["info"]["pdf_version"] = header.decode('latin-1')
    except Exception as e:
        results["errors"].append(f"Cannot read file: {e}")
        return results
    
    # Check 3: PDF EOF marker
    try:
        with open(filepath, 'rb') as f:
            f.seek(-1024, 2)  # Last 1KB
            tail = f.read()
            if b'%%EOF' not in tail:
                results["warnings"].append("No EOF marker found (file may be truncated)")
    except Exception as e:
        results["warnings"].append(f"Cannot check EOF: {e}")
    
    # Mark as valid if no errors
    results["valid"] = len(results["errors"]) == 0
    
    return results


def find_source_templates():
    """Find source template files."""
    templates = []
    templates_dir = os.path.join(ROOT, "assets", "slides", "templates")
    
    if os.path.exists(templates_dir):
        for filename in sorted(os.listdir(templates_dir)):
            if filename.lower().endswith('.pptx') and not filename.startswith('~$'):
                templates.append(os.path.join(templates_dir, filename))
    
    return templates


def validate_directory(directory=None, check_structure=True):
    """Validate all output files in a directory.
    
    Args:
        directory: Directory to validate (default: OUTPUT_DIR)
        check_structure: Whether to compare against source templates
    """
    if directory is None:
        directory = OUTPUT_DIR
    
    if not os.path.exists(directory):
        print(f"❌ Directory does not exist: {directory}")
        return []
    
    results = []
    source_templates = find_source_templates() if check_structure else None
    
    # Walk through directory tree
    for root, dirs, files in os.walk(directory):
        for filename in sorted(files):
            if filename.startswith('.') or filename.startswith('~$'):
                continue
            
            filepath = os.path.join(root, filename)
            
            if filename.lower().endswith('.pptx'):
                result = validate_pptx(filepath, source_templates)
                results.append(result)
            elif filename.lower().endswith('.pdf'):
                result = validate_pdf(filepath)
                results.append(result)
    
    return results


def print_results(results):
    """Pretty print validation results."""
    if not results:
        print("No files found to validate.")
        return
    
    valid_count = sum(1 for r in results if r["valid"])
    total_count = len(results)
    
    print(f"\n{'='*70}")
    print(f"VALIDATION REPORT: {valid_count}/{total_count} files valid")
    print(f"{'='*70}\n")
    
    for result in results:
        status = "✓" if result["valid"] else "✗"
        print(f"{status} {result['file']}")
        
        if result["info"]:
            for key, value in result["info"].items():
                print(f"    {key}: {value}")
        
        if result["errors"]:
            for error in result["errors"]:
                print(f"    ❌ ERROR: {error}")
        
        if result["warnings"]:
            for warning in result["warnings"]:
                print(f"    ⚠️  WARNING: {warning}")
        
        print()
    
    # Summary
    if valid_count == total_count:
        print("✅ All files are valid!")
        return 0
    else:
        print(f"⚠️  {total_count - valid_count} file(s) have issues")
        return 1


def main():
    import argparse
    parser = argparse.ArgumentParser(description="Validate merged output files with structure checks")
    parser.add_argument("path", nargs="?", help="File or directory to validate (default: assets/output)")
    parser.add_argument("--quiet", "-q", action="store_true", help="Only show errors")
    parser.add_argument("--no-structure", action="store_true", help="Skip structure comparison with templates")
    args = parser.parse_args()
    
    target = args.path or OUTPUT_DIR
    check_structure = not args.no_structure
    
    if os.path.isfile(target):
        # Validate single file
        source_templates = find_source_templates() if check_structure else None
        
        if target.lower().endswith('.pptx'):
            result = validate_pptx(target, source_templates)
        elif target.lower().endswith('.pdf'):
            result = validate_pdf(target)
        else:
            print(f"❌ Unsupported file type: {target}")
            return 1
        
        results = [result]
    else:
        # Validate directory
        results = validate_directory(target, check_structure)
    
    return print_results(results)


if __name__ == "__main__":
    sys.exit(main())
