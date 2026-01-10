#!/usr/bin/env python3
"""Inspect a PPTX and print helpful diagnostics for debugging merge issues.

Usage: python3 src/inspect_pptx.py path/to/file.pptx
"""
import sys
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def inspect(path):
    print(f"\nInspecting: {path}")
    if not os.path.isfile(path):
        print("  File not found")
        return
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"  Failed to open presentation: {e!r}")
        return

    # Basic counts
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Slide masters: {len(masters)}")
    try:
        layouts = len(prs.slide_layouts)
    except Exception:
        layouts = 0
    print(f"  Slide layouts: {layouts}")

    # Core properties
    try:
        cp = prs.core_properties
        props = {k: getattr(cp, k) for k in ("author", "title", "subject", "keywords")}
        print(f"  Core properties: {props}")
    except Exception:
        print("  Could not read core properties")

    # Count package parts (best-effort)
    try:
        pkg_parts = len(prs.part.package.parts)
        print(f"  Package parts: {pkg_parts}")
    except Exception:
        pass

    # Inspect each slide and shapes
    for i, slide in enumerate(prs.slides, start=1):
        print(f"\n  Slide {i}: shapes={len(slide.shapes)}")
        for j, shape in enumerate(slide.shapes, start=1):
            stype = getattr(shape, "shape_type", None)
            info = f"    shape {j}: type={stype}"
            try:
                if stype == MSO_SHAPE_TYPE.PICTURE:
                    # picture
                    try:
                        img = shape.image
                        blob_size = len(img.blob)
                        info += f", picture (size={blob_size} bytes, ext={img.ext})"
                    except Exception as e:
                        info += f", picture (could not read blob: {e!r})"
                elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text
                    info += f", text_len={len(text)}"
                elif getattr(shape, "has_chart", False):
                    info += ", chart"
                else:
                    # other types: graphicFrame, table, group, etc.
                    info += f", other (has_table={getattr(shape, 'has_table', False)}, has_chart={getattr(shape, 'has_chart', False)})"
            except Exception as e:
                info += f", error_inspecting={e!r}"
            print(info)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 src/inspect_pptx.py file.pptx")
        sys.exit(1)
    for p in sys.argv[1:]:
        inspect(p)
