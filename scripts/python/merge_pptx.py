#!/usr/bin/env python3
"""Merge multiple PPTX files into one, saving into `assets/output`.

Usage:
  python merge_pptx.py [output_name.pptx] [deck1.pptx deck2.pptx ...]
  python merge_pptx.py --config path/to/config.json

If no output name is provided, generates: merged-pptx-{number}.pptx
If no input decks are provided, all files in `assets/slides/templates` will be used.
If --config is provided, reads merge configuration from JSON file.
The resulting file will be written to `assets/output/`.
"""
import argparse
import io
import os
import sys
import json
import random
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# Import theme resolver for explicit formatting
sys.path.insert(0, os.path.join(os.path.dirname(__file__)))
from theme_resolver import apply_explicit_formatting


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
ASSETS_DIR = os.path.join(ROOT, "assets")
ASSETS_SLIDES = os.path.join(ROOT, "assets", "slides")
TEMPLATES_DIR = os.path.join(ASSETS_SLIDES, "templates")
CONFIG_DIR = os.path.join(ROOT, "assets", "config")
OUTPUT_DIR = os.path.join(ROOT, "assets", "output", "merge_pptx")


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def generate_output_filename(prefix="merged-pptx", extension=".pptx"):
    """Generate incremental filename in assets/output."""
    ensure_dir(OUTPUT_DIR)
    counter = 1
    while True:
        filename = f"{prefix}-{counter:03d}{extension}"
        filepath = os.path.join(OUTPUT_DIR, filename)
        if not os.path.exists(filepath):
            return filepath
        counter += 1
        if counter > 999:
            # Fallback to random
            filename = f"{prefix}-{random.randint(1000, 9999)}{extension}"
            return os.path.join(OUTPUT_DIR, filename)


def find_template():
    if not os.path.isdir(TEMPLATES_DIR):
        return None
    for name in os.listdir(TEMPLATES_DIR):
        if name.lower().endswith(".pptx"):
            return os.path.join(TEMPLATES_DIR, name)
    return None


def copy_text(shape, target_slide):
    # Keep as a fallback. Prefer deep-copying the shape XML to preserve
    # formatting when possible (see append_slide_from_source).
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        textbox = target_slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        # Simple copy: join paragraphs with line breaks
        parts = []
        for p in shape.text_frame.paragraphs:
            runs_text = "".join([r.text for r in p.runs])
            parts.append(runs_text)
        tf.text = "\n".join(parts)
    except Exception:
        pass


def copy_picture(shape, target_slide):
    try:
        image = shape.image
        blob = image.blob
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        target_slide.shapes.add_picture(io.BytesIO(blob), left, top, width, height)
    except Exception:
        pass


def copy_shape(shape, target_slide):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        copy_picture(shape, target_slide)
    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
        copy_text(shape, target_slide)
    else:
        # Not handling groups, charts, tables robustly.
        return


def get_master_font_size(source_slide, placeholder_type='body', level=1):
    """Get the default font size from the source slide's master for a given placeholder type."""
    try:
        master = source_slide.slide_layout.slide_master
        master_elem = master._element
        ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        
        txStyles = master_elem.find(f'{ns}txStyles')
        if txStyles is not None:
            # Choose the right style based on placeholder type
            if placeholder_type == 'title':
                style = txStyles.find(f'{ns}titleStyle')
            elif placeholder_type == 'body':
                style = txStyles.find(f'{ns}bodyStyle')
            else:
                style = txStyles.find(f'{ns}otherStyle')
            
            if style is not None:
                # Get the level (lvl1pPr, lvl2pPr, etc.)
                lvlpPr = style.find(f'{ns_a}lvl{level}pPr')
                if lvlpPr is not None:
                    defRPr = lvlpPr.find(f'{ns_a}defRPr')
                    if defRPr is not None:
                        sz = defRPr.get('sz')
                        if sz:
                            # Size is in hundredths of a point
                            return int(sz) / 100
    except:
        pass
    return None


def get_theme_font_name(source_slide, is_major=False):
    """Get the theme font name (major for headings, minor for body) from source slide's master."""
    try:
        master = source_slide.slide_layout.slide_master
        master_elem = master._element
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        
        # Find themeElements
        themeElements = master_elem.find(f'.//{ns_a}themeElements')
        if themeElements is not None:
            fontScheme = themeElements.find(f'{ns_a}fontScheme')
            if fontScheme is not None:
                if is_major:
                    font = fontScheme.find(f'{ns_a}majorFont')
                else:
                    font = fontScheme.find(f'{ns_a}minorFont')
                
                if font is not None:
                    latin = font.find(f'{ns_a}latin')
                    if latin is not None:
                        return latin.get('typeface')
    except:
        pass
    return None


def append_slide_from_source(merged_presentation, source_slide, source_presentation):
    # Create a new slide and copy shapes properly, preserving formatting
    # and handling image relationships correctly
    
    # Try to find a matching layout by name from source, otherwise use blank
    layout = None
    source_layout_name = source_slide.slide_layout.name
    
    # First try to find matching layout by name
    for candidate_layout in merged_presentation.slide_layouts:
        if candidate_layout.name == source_layout_name:
            layout = candidate_layout
            break
    
    # If no match, look for a blank layout
    if layout is None:
        for candidate_layout in merged_presentation.slide_layouts:
            if 'blank' in candidate_layout.name.lower():
                layout = candidate_layout
                break
    
    # Fallback to layout 6 (commonly blank) or layout 0
    if layout is None:
        try:
            layout = merged_presentation.slide_layouts[6] if len(merged_presentation.slide_layouts) > 6 else merged_presentation.slide_layouts[0]
        except:
            layout = merged_presentation.slide_layouts[0]

    target_slide = merged_presentation.slides.add_slide(layout)
    
    # Copy slide background from source - preserve the exact background
    try:
        # Clone the entire background element from source slide
        source_cSld = source_slide.element.cSld
        target_cSld = target_slide.element.cSld
        
        # Check if source has background properties
        if source_cSld.bg is not None:
            # Remove existing background if present
            if target_cSld.bg is not None:
                target_cSld.remove(target_cSld.bg)
            # Clone and add source background
            target_cSld.insert(0, deepcopy(source_cSld.bg))
        else:
            # Source doesn't override background (uses master)
            # We need to copy the actual rendered background color
            # to ensure it looks the same even if masters differ
            try:
                source_fill = source_slide.background.fill
                if source_fill.type:
                    target_slide.background.fill.solid()
                    # Try to get the actual color
                    if source_fill.type == 1:  # SOLID
                        if source_fill.fore_color.type == 1:  # RGB
                            target_slide.background.fill.fore_color.rgb = source_fill.fore_color.rgb
                        elif source_fill.fore_color.type == 2:  # THEME/SCHEME
                            # Copy theme color reference
                            target_slide.background.fill.fore_color.theme_color = source_fill.fore_color.theme_color
                    elif source_fill.type == 5:  # BACKGROUND (inherits from master)
                        # Need to get the actual resolved color from the master
                        # For now, try to access master background
                        try:
                            master_fill = source_slide.slide_layout.slide_master.background.fill
                            if master_fill.type == 1:  # SOLID
                                target_slide.background.fill.solid()
                                if master_fill.fore_color.type == 1:
                                    target_slide.background.fill.fore_color.rgb = master_fill.fore_color.rgb
                                elif master_fill.fore_color.type == 2:
                                    target_slide.background.fill.fore_color.theme_color = master_fill.fore_color.theme_color
                        except:
                            pass
            except Exception as e2:
                print(f"  ⚠️  Could not copy master background color: {e2}")
    except Exception as e:
        print(f"  ⚠️  Could not copy slide background: {e}")
    
    # Remove any placeholder shapes from the new slide to start clean
    # This prevents duplicates when we copy shapes from the source
    shapes_to_remove = [sp for sp in target_slide.shapes]
    for shape in shapes_to_remove:
        sp = shape.element
        sp.getparent().remove(sp)

    # Copy each shape - for pictures, we need special handling to preserve the image
    for shape in source_slide.shapes:
        try:
            # Check if it's a picture - these need special handling for image relationships
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Copy picture with its image data
                image = shape.image
                blob = image.blob
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                new_pic = target_slide.shapes.add_picture(io.BytesIO(blob), left, top, width, height)
                # Try to copy the name
                try:
                    new_pic.name = shape.name
                except:
                    pass
            else:
                # For non-picture shapes, try deepcopy to preserve formatting
                el = deepcopy(shape._element)
                target_slide.shapes._spTree.append(el)
                
                # CRITICAL: Preserve bullet formatting from source
                # If the source paragraph doesn't have explicit bullet settings,
                # we need to check the source and copy that state
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    # Find the corresponding shape in target
                    target_shape = None
                    for ts in target_slide.shapes:
                        if ts.name == shape.name:
                            target_shape = ts
                            break
                    
                    if target_shape and hasattr(target_shape, 'has_text_frame') and target_shape.has_text_frame:
                        # Copy bullet settings and font properties from each paragraph
                        for src_para, tgt_para in zip(shape.text_frame.paragraphs, target_shape.text_frame.paragraphs):
                            src_pPr = src_para._element.pPr
                            tgt_pPr = tgt_para._element.get_or_add_pPr()
                            
                            # Check if source has any bullet element
                            has_bullet_element = False
                            if src_pPr is not None:
                                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                                buNone = src_pPr.find(f'.//{ns}buNone')
                                buChar = src_pPr.find(f'.//{ns}buChar')
                                buAutoNum = src_pPr.find(f'.//{ns}buAutoNum')
                                has_bullet_element = any([buNone is not None, buChar is not None, buAutoNum is not None])
                            
                            # If source has NO explicit bullet element and it's level 0,
                            # it means "no bullets" - add buNone to target
                            if not has_bullet_element and src_para.level == 0:
                                # Remove any existing bullet elements from target
                                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                                for bullet_elem in tgt_pPr.findall(f'.//{ns}buChar') + tgt_pPr.findall(f'.//{ns}buAutoNum'):
                                    tgt_pPr.remove(bullet_elem)
                                # Add buNone if not already there
                                if tgt_pPr.find(f'.//{ns}buNone') is None:
                                    from lxml import etree
                                    buNone = etree.SubElement(tgt_pPr, f'{ns}buNone')
                            
                            # CRITICAL: Apply explicit formatting from source
                            # Resolve all theme-dependent values (colors, fonts, sizes) to explicit values
                            # This ensures exact appearance preservation regardless of theme differences
                            for src_run, tgt_run in zip(src_para.runs, tgt_para.runs):
                                apply_explicit_formatting(
                                    src_run, tgt_run, 
                                    source_presentation, source_slide,
                                    shape, src_para
                                )

        except Exception as e:
            # Deepcopy failed, try manual copy
            try:
                copy_shape(shape, target_slide)
            except Exception as e2:
                # Log but continue - don't let one shape failure stop the whole merge
                print(f"  ⚠️  Could not copy shape {shape.name if hasattr(shape, 'name') else 'unnamed'}: {e2}")


def load_config(config_path):
    """Load merge configuration from JSON file.
    
    Expected format:
    {
      "output": "merged-output.pptx",
      "parts": [
        {
          "file": "slides/templates/test1.pptx",
          "slides": "all" or [0, 1, 2]
        },
        {
          "file": "slides/templates/test2.pptx",
          "slides": "all"
        }
      ]
    }
    
    Or simplified:
    {
      "output": "merged-output.pptx",
      "parts": [
        "slides/templates/test1.pptx",
        "slides/templates/test2.pptx"
      ]
    }
    """
    with open(config_path, 'r') as f:
        config = json.load(f)
    
    # Normalize parts to dict format
    normalized_parts = []
    for part in config.get('parts', []):
        if isinstance(part, str):
            # Simple string format
            normalized_parts.append({
                'file': part,
                'slides': 'all'
            })
        else:
            # Dict format
            normalized_parts.append(part)
    
    # Resolve relative paths to assets directory
    for part in normalized_parts:
        file_path = part['file']
        # Try multiple resolution strategies
        if not os.path.isabs(file_path):
            # Try relative to assets
            candidate = os.path.join(ASSETS_DIR, file_path)
            if os.path.exists(candidate):
                part['file'] = candidate
            else:
                # Try relative to root
                candidate = os.path.join(ROOT, file_path)
                if os.path.exists(candidate):
                    part['file'] = candidate
    
    return config.get('output'), normalized_parts


def merge_presentations(output_file, input_files, parts_config=None):
    # Use the first input file as the template to inherit slide masters and layouts
    # but start with a blank presentation and copy ALL slides properly
    template_path = find_template()
    
    if input_files or parts_config:
        # Determine which files to process
        files_to_process = []
        
        if parts_config:
            # Use config-based processing
            for part in parts_config:
                file_path = part['file']
                slides_spec = part.get('slides', 'all')
                files_to_process.append((file_path, slides_spec))
        else:
            # Use simple file list
            files_to_process = [(f, 'all') for f in input_files]
        
        if not files_to_process:
            print("Warning: No files to process")
            return
        
        # Load first file to get the template structure (masters, layouts, theme)
        first_file = files_to_process[0][0]
        base_prs = Presentation(first_file)
        # Create merged presentation inheriting from the first file's structure
        merged = Presentation(first_file)
        # Remove all slides from the base - we'll add them back properly
        while len(merged.slides) > 0:
            rId = merged.slides._sldIdLst[0].rId
            merged.part.drop_rel(rId)
            del merged.slides._sldIdLst[0]
        
        # Now copy slides from all input files according to config
        for file_path, slides_spec in files_to_process:
            if not os.path.isfile(file_path):
                print(f"Warning: input file not found: {file_path}")
                continue
            
            src = Presentation(file_path)
            
            if slides_spec == 'all':
                # Copy all slides
                for slide in src.slides:
                    append_slide_from_source(merged, slide, src)
            else:
                # Copy specific slides by index
                for slide_idx in slides_spec:
                    if 0 <= slide_idx < len(src.slides):
                        append_slide_from_source(merged, src.slides[slide_idx], src)
                    else:
                        print(f"Warning: slide index {slide_idx} out of range in {file_path}")
    elif template_path:
        merged = Presentation(template_path)
    else:
        merged = Presentation()
        print("Warning: No input files provided and no template found")

    ensure_dir(os.path.dirname(output_file))
    merged.save(output_file)
    print(f"Saved merged presentation to: {output_file}")
    
    # Validate the output
    try:
        test = Presentation(output_file)
        print(f"✓ Validation: {len(test.slides)} slides, {os.path.getsize(output_file)} bytes")
    except Exception as e:
        print(f"⚠️  Validation warning: {e}")


def main():
    p = argparse.ArgumentParser(description="Merge PPTX files into one in assets/output")
    p.add_argument("output", nargs="?", help="Output PPTX filename (saved under assets/output). If omitted, auto-generates name.")
    p.add_argument("inputs", nargs="*", help="Input PPTX files to merge. If omitted, templates are used.")
    p.add_argument("--config", "-c", help="JSON config file specifying merge order and options")
    p.add_argument("--template", "-t", help="Optional template/master PPTX to use as the base for the merged presentation")
    args = p.parse_args()

    ensure_dir(OUTPUT_DIR)

    # Check if config file is provided
    parts_config = None
    config_output = None
    
    if args.config:
        config_path = args.config
        # Try to resolve config path relative to config dir if not absolute
        if not os.path.isabs(config_path):
            candidate = os.path.join(CONFIG_DIR, config_path)
            if os.path.exists(candidate):
                config_path = candidate
        
        if not os.path.exists(config_path):
            print(f"Config file not found: {config_path}")
            sys.exit(1)
        
        config_output, parts_config = load_config(config_path)
        print(f"Loaded config from: {config_path}")
        print(f"  Parts: {len(parts_config)} file(s)")

    # Resolve input files
    input_files = args.inputs or []
    # If no explicit inputs were provided and no config, gather templates from the templates dir
    if not input_files and not parts_config:
        if os.path.isdir(TEMPLATES_DIR):
            # Filter out PowerPoint lock/temp files (starting with ~$)
            input_files = [
                os.path.join(TEMPLATES_DIR, n) 
                for n in sorted(os.listdir(TEMPLATES_DIR))
                if n.lower().endswith(".pptx") and not n.startswith("~$")
            ]

    # If a template was explicitly provided via CLI, set it and ensure it exists
    cli_template = None
    if args.template:
        if os.path.isfile(args.template):
            cli_template = args.template
        else:
            print(f"Template file not found: {args.template}")
            sys.exit(1)

    # Check if we have any input
    if not input_files and not parts_config:
        print("No input presentations found. Provide files, use --config, or add templates to assets/slides/templates.")
        sys.exit(1)

    # Determine the output filename
    if config_output:
        # Config specified output name
        output_filename = config_output
    elif args.output:
        output_filename = args.output
    else:
        # Auto-generate filename
        output_filename = os.path.basename(generate_output_filename())

    # If output doesn't end with .pptx, add it
    if not output_filename.lower().endswith(".pptx"):
        output_filename += ".pptx"

    # Ensure output is in OUTPUT_DIR
    if os.path.isabs(output_filename):
        output_path = output_filename
    else:
        output_path = os.path.join(OUTPUT_DIR, output_filename)

    # If the user passed a template, prepend it to input_files
    if cli_template and input_files:
        input_files = [cli_template] + input_files
    
    # Run the merge
    merge_presentations(output_path, input_files, parts_config)


if __name__ == "__main__":
    main()
