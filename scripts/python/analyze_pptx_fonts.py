#!/usr/bin/env python3
"""
Analyze PPTX files for font attribute outliers.
Reads files without modifying them and generates analysis reports.

Usage:
  python analyze_pptx_fonts.py <source_dir> [--since YYYY-MM-DD] [--output output_dir]
"""
import argparse
import os
import sys
import json
from datetime import datetime
from collections import defaultdict, Counter
from pathlib import Path
from pptx import Presentation

# Add parent directory to path
sys.path.insert(0, os.path.dirname(__file__))
from theme_resolver import get_theme_color_rgb, get_master_font_size


def get_file_mtime(filepath):
    """Get file modification time."""
    return datetime.fromtimestamp(os.path.getmtime(filepath))


def analyze_text_run(run, file_path, slide_idx, shape_idx, para_idx, run_idx, presentation, slide, shape):
    """Extract all font attributes from a text run with deep theme resolution."""
    attrs = {
        'file': os.path.basename(file_path),
        'file_path': file_path,
        'slide': slide_idx,
        'shape': shape_idx,
        'paragraph': para_idx,
        'run': run_idx,
        'text': run.text[:50] if run.text else '',  # First 50 chars
        'text_length': len(run.text) if run.text else 0,
    }
    
    # Font attributes
    try:
        # Font name (surface value)
        attrs['font_name'] = run.font.name if run.font.name else '(theme)'
        
        # Font size with master resolution
        if run.font.size:
            attrs['font_size'] = run.font.size.pt
            attrs['font_size_source'] = 'explicit'
        else:
            # Try to get from master
            try:
                placeholder_type = 'body'
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    if ph_type == 1:  # TITLE
                        placeholder_type = 'title'
                
                master_size = get_master_font_size(slide, placeholder_type, 1)
                if master_size:
                    attrs['font_size'] = master_size
                    attrs['font_size_source'] = 'master'
                else:
                    attrs['font_size'] = None
                    attrs['font_size_source'] = 'none'
            except:
                attrs['font_size'] = None
                attrs['font_size_source'] = 'none'
        
        attrs['bold'] = run.font.bold if run.font.bold is not None else None
        attrs['italic'] = run.font.italic if run.font.italic is not None else None
        attrs['underline'] = run.font.underline if run.font.underline is not None else None
        
        # Color with theme resolution
        if run.font.color.type == 1:  # RGB
            attrs['color_type'] = 'RGB'
            attrs['color_value'] = str(run.font.color.rgb)
            attrs['color_resolved'] = str(run.font.color.rgb)
        elif run.font.color.type == 2:  # SCHEME
            attrs['color_type'] = 'SCHEME'
            attrs['color_value'] = f'SCHEME({run.font.color.theme_color})'
            # Resolve to actual RGB
            resolved_rgb = get_theme_color_rgb(presentation, run.font.color.theme_color)
            attrs['color_resolved'] = resolved_rgb if resolved_rgb else attrs['color_value']
        else:
            attrs['color_type'] = 'NONE'
            attrs['color_value'] = None
            attrs['color_resolved'] = None
    except Exception as e:
        attrs['error'] = str(e)
    
    return attrs


def analyze_pptx_file(filepath):
    """Analyze a single PPTX file."""
    try:
        prs = Presentation(filepath)
        runs_data = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        for run_idx, run in enumerate(para.runs):
                            if run.text.strip():  # Only analyze non-empty runs
                                run_data = analyze_text_run(
                                    run, filepath, slide_idx, shape_idx, para_idx, run_idx,
                                    prs, slide, shape
                                )
                                runs_data.append(run_data)
        
        return {
            'success': True,
            'file': filepath,
            'slides': len(prs.slides),
            'runs': runs_data
        }
    except Exception as e:
        return {
            'success': False,
            'file': filepath,
            'error': str(e)
        }


def find_outliers(all_runs):
    """Identify outliers in font attributes."""
    # Collect statistics
    font_sizes = [r['font_size'] for r in all_runs if r.get('font_size')]
    font_names = [r['font_name'] for r in all_runs if r.get('font_name')]
    # Use resolved colors (SCHEME colors are resolved to RGB)
    colors = [r['color_resolved'] for r in all_runs if r.get('color_resolved')]
    
    # Count occurrences
    size_counts = Counter(font_sizes)
    name_counts = Counter(font_names)
    color_counts = Counter(colors)
    
    # Determine common values (>10% of total)
    total_runs = len(all_runs)
    threshold = total_runs * 0.10
    
    common_sizes = {size for size, count in size_counts.items() if count > threshold}
    common_names = {name for name, count in name_counts.items() if count > threshold}
    common_colors = {color for color, count in color_counts.items() if count > threshold}
    
    # Find outliers
    outliers = {
        'size': [],
        'name': [],
        'color': []
    }
    
    for run in all_runs:
        # Size outliers
        if run.get('font_size') and run['font_size'] not in common_sizes:
            outliers['size'].append(run)
        
        # Font name outliers
        if run.get('font_name') and run['font_name'] not in common_names:
            outliers['name'].append(run)
        
        # Color outliers (use resolved values)
        if run.get('color_resolved') and run['color_resolved'] not in common_colors:
            outliers['color'].append(run)
    
    return {
        'statistics': {
            'total_runs': total_runs,
            'font_sizes': dict(size_counts.most_common(20)),
            'font_names': dict(name_counts.most_common(20)),
            'colors': dict(color_counts.most_common(20)),
            'common_sizes': list(common_sizes),
            'common_names': list(common_names),
            'common_colors': list(common_colors)
        },
        'outliers': outliers
    }


def generate_report(analysis_results, outlier_analysis, output_dir):
    """Generate analysis reports."""
    os.makedirs(output_dir, exist_ok=True)
    
    # Save full data
    with open(os.path.join(output_dir, 'full_analysis.json'), 'w') as f:
        json.dump({
            'files_analyzed': len(analysis_results),
            'results': analysis_results
        }, f, indent=2)
    
    # Save outlier analysis
    with open(os.path.join(output_dir, 'outlier_analysis.json'), 'w') as f:
        json.dump(outlier_analysis, f, indent=2)
    
    # Generate summary report
    summary_lines = []
    summary_lines.append("=" * 80)
    summary_lines.append("PPTX FONT ATTRIBUTE ANALYSIS REPORT")
    summary_lines.append("=" * 80)
    summary_lines.append(f"Files analyzed: {len(analysis_results)}")
    
    successful = [r for r in analysis_results if r['success']]
    failed = [r for r in analysis_results if not r['success']]
    
    summary_lines.append(f"Successful: {len(successful)}")
    summary_lines.append(f"Failed: {len(failed)}")
    
    if failed:
        summary_lines.append("\nFailed files:")
        for f in failed[:10]:
            summary_lines.append(f"  - {f['file']}: {f.get('error', 'Unknown error')}")
    
    stats = outlier_analysis['statistics']
    summary_lines.append(f"\nTotal text runs analyzed: {stats['total_runs']}")
    
    summary_lines.append("\n" + "=" * 80)
    summary_lines.append("FONT SIZE DISTRIBUTION")
    summary_lines.append("=" * 80)
    for size, count in list(stats['font_sizes'].items())[:10]:
        pct = (count / stats['total_runs']) * 100
        summary_lines.append(f"  {size}pt: {count} ({pct:.1f}%)")
    
    summary_lines.append("\n" + "=" * 80)
    summary_lines.append("FONT FAMILY DISTRIBUTION")
    summary_lines.append("=" * 80)
    for name, count in list(stats['font_names'].items())[:10]:
        pct = (count / stats['total_runs']) * 100
        summary_lines.append(f"  {name}: {count} ({pct:.1f}%)")
    
    summary_lines.append("\n" + "=" * 80)
    summary_lines.append("COLOR DISTRIBUTION")
    summary_lines.append("=" * 80)
    for color, count in list(stats['colors'].items())[:10]:
        pct = (count / stats['total_runs']) * 100
        summary_lines.append(f"  {color}: {count} ({pct:.1f}%)")
    
    summary_lines.append("\n" + "=" * 80)
    summary_lines.append("OUTLIERS DETECTED")
    summary_lines.append("=" * 80)
    summary_lines.append(f"Font size outliers: {len(outlier_analysis['outliers']['size'])}")
    summary_lines.append(f"Font name outliers: {len(outlier_analysis['outliers']['name'])}")
    summary_lines.append(f"Color outliers: {len(outlier_analysis['outliers']['color'])}")
    
    # Show sample outliers
    if outlier_analysis['outliers']['size']:
        summary_lines.append("\nSample size outliers:")
        for run in outlier_analysis['outliers']['size'][:5]:
            summary_lines.append(f"  {run['file']} (slide {run['slide']}): {run['font_size']}pt - \"{run['text']}\"")
    
    if outlier_analysis['outliers']['name']:
        summary_lines.append("\nSample font name outliers:")
        for run in outlier_analysis['outliers']['name'][:5]:
            summary_lines.append(f"  {run['file']} (slide {run['slide']}): {run['font_name']} - \"{run['text']}\"")
    
    summary_lines.append("\n" + "=" * 80)
    summary_lines.append(f"Full results saved to: {output_dir}")
    summary_lines.append("=" * 80)
    
    summary_text = "\n".join(summary_lines)
    
    # Save summary
    with open(os.path.join(output_dir, 'summary.txt'), 'w') as f:
        f.write(summary_text)
    
    # Print summary
    print(summary_text)


def main():
    parser = argparse.ArgumentParser(description="Analyze PPTX files for font attribute outliers")
    parser.add_argument('source_dir', help="Directory containing PPTX files")
    parser.add_argument('--since', help="Only analyze files modified since this date (YYYY-MM-DD)")
    parser.add_argument('--output', help="Output directory for analysis results", 
                        default='assets/analysis')
    args = parser.parse_args()
    
    # Parse date filter
    since_date = None
    if args.since:
        since_date = datetime.strptime(args.since, '%Y-%m-%d')
    
    # Find PPTX files
    pptx_files = []
    for root, dirs, files in os.walk(args.source_dir):
        for file in files:
            if file.endswith('.pptx') and not file.startswith('~$'):
                filepath = os.path.join(root, file)
                if since_date:
                    if get_file_mtime(filepath) >= since_date:
                        pptx_files.append(filepath)
                else:
                    pptx_files.append(filepath)
    
    print(f"Found {len(pptx_files)} PPTX files to analyze")
    
    if not pptx_files:
        print("No files found matching criteria")
        return
    
    # Analyze files
    print("Analyzing files...")
    analysis_results = []
    all_runs = []
    
    for i, filepath in enumerate(pptx_files, 1):
        print(f"  [{i}/{len(pptx_files)}] {os.path.basename(filepath)}")
        result = analyze_pptx_file(filepath)
        analysis_results.append(result)
        
        if result['success']:
            all_runs.extend(result['runs'])
    
    print(f"\nAnalyzed {len(all_runs)} text runs across {len(pptx_files)} files")
    
    # Find outliers
    print("Identifying outliers...")
    outlier_analysis = find_outliers(all_runs)
    
    # Generate report
    print("Generating reports...")
    generate_report(analysis_results, outlier_analysis, args.output)


if __name__ == "__main__":
    main()
