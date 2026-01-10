#!/usr/bin/env python3
"""Create a PPTX from a directory of images (PNG/JPG), one slide per image.

Usage: python3 src/pdf_images_to_pptx.py output.pptx images_dir
"""
import sys
import os
from pptx import Presentation


def images_from_dir(images_dir):
    files = sorted([f for f in os.listdir(images_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
    return [os.path.join(images_dir, f) for f in files]


def build_pptx(output_path, images):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for img in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)


def main():
    if len(sys.argv) != 3:
        print("Usage: python3 src/pdf_images_to_pptx.py output.pptx images_dir")
        sys.exit(2)
    out = sys.argv[1]
    images_dir = sys.argv[2]
    if not os.path.isdir(images_dir):
        print(f"Images directory not found: {images_dir}")
        sys.exit(1)
    images = images_from_dir(images_dir)
    if not images:
        print("No images found in directory")
        sys.exit(1)
    build_pptx(out, images)


if __name__ == '__main__':
    main()
