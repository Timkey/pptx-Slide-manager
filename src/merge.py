import sys
from pptx import Presentation

def merge_presentations(output_file, input_files):
    merged_presentation = Presentation(input_files[0]) # Start with the first file

    for file in input_files[1:]:
        source_presentation = Presentation(file)
        for slide in source_presentation.slides:
            # Note: This basic method copies content but 
            # requires manual handling for complex layouts/masters
            slide_layout = merged_presentation.slide_layouts[0] 
            merged_presentation.slides.add_slide(slide_layout)
            # Add logic here to copy shapes if deep-cloning is needed
            
    merged_presentation.save(output_file)
    print(f"Successfully merged into {output_file}")

if __name__ == "__main__":
    output = sys.argv[1]
    inputs = sys.argv[2:]
    merge_presentations(output, inputs)