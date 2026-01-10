#!/usr/bin/env python3
"""Merge multiple PPTX files into one, saving into `assets/slides/merged`.

Usage:
  python merge_pptx.py final_deck.pptx deck1.pptx deck2.pptx ...

If no input decks are provided, all files in `assets/slides/templates` will be used.
The resulting file will be written to `assets/slides/merged/<output>`.
"""
import argparse
import io
import os
import sys
from copy import deepcopy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
ASSETS_SLIDES = os.path.join(ROOT, "assets", "slides")
TEMPLATES_DIR = os.path.join(ASSETS_SLIDES, "templates")
MERGED_DIR = os.path.join(ASSETS_SLIDES, "merged")


def ensure_dir(path):
    os.makedirs(path, exist_ok=True)


def find_template():
    if not os.path.isdir(TEMPLATES_DIR):
        return None
    for name in os.listdir(TEMPLATES_DIR):
        if name.lower().endswith(".pptx"):
            return os.path.join(TEMPLATES_DIR, name)
    return None


def copy_text(shape, target_slide):
    # Keep as a fallback. Prefer deep-copying the shape XML to preserve
    # formatting when possible (see append_slide_from_source).
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        textbox = target_slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        # Simple copy: join paragraphs with line breaks
        parts = []
        for p in shape.text_frame.paragraphs:
            runs_text = "".join([r.text for r in p.runs])
            parts.append(runs_text)
        tf.text = "\n".join(parts)
    except Exception:
        pass


def copy_picture(shape, target_slide):
    try:
        image = shape.image
        blob = image.blob
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        target_slide.shapes.add_picture(io.BytesIO(blob), left, top, width, height)
    except Exception:
        pass


def copy_shape(shape, target_slide):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        copy_picture(shape, target_slide)
    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
        copy_text(shape, target_slide)
    else:
        # Not handling groups, charts, tables robustly.
        return


def append_slide_from_source(merged_presentation, source_slide):
    # Create a new slide and deep-copy the shapes' XML elements. This
    # preserves most visual formatting (fonts, fills, etc.). Complex
    # objects (charts, embedded OLE, advanced SmartArt) may still need
    # special handling.
    try:
        layout = merged_presentation.slide_layouts[6] if len(merged_presentation.slide_layouts) > 6 else merged_presentation.slide_layouts[0]
    except Exception:
        layout = merged_presentation.slide_layouts[0]

    target_slide = merged_presentation.slides.add_slide(layout)

    # Deep-copy each shape element into the target slide's spTree. If that
    # fails for a shape, fall back to manual copying.
    spTree = target_slide.shapes._spTree
    for shape in source_slide.shapes:
        try:
            el = deepcopy(shape._element)
            spTree.append(el)
        except Exception:
            copy_shape(shape, target_slide)


def merge_presentations(output_file, input_files):
    # Determine a base presentation for the merged output.
    # If explicit input files are provided, use the first input as the base
    # to avoid duplicating a template when templates are also passed as inputs.
    template_path = find_template()
    # Note: caller may provide a template via CLI (handled in main), so here
    # we only look at input_files and templates directory as fallback.
    base_from_input = False
    if input_files:
        merged = Presentation(input_files[0])
        base_from_input = True
    elif template_path:
        merged = Presentation(template_path)
    else:
        merged = Presentation()
    # If merged was initialized from the first input, skip that file when
    # copying slides to avoid duplication. Otherwise copy from the first.
    start_idx = 1 if base_from_input else 0
    for infile in input_files[start_idx:]:
        if not os.path.isfile(infile):
            print(f"Warning: input file not found: {infile}")
            continue
        src = Presentation(infile)
        for slide in src.slides:
            append_slide_from_source(merged, slide)

    ensure_dir(os.path.dirname(output_file))
    merged.save(output_file)
    print(f"Saved merged presentation to: {output_file}")


def main():
    p = argparse.ArgumentParser(description="Merge PPTX files into one in assets/slides/merged")
    p.add_argument("output", help="Output PPTX filename (saved under assets/slides/merged)")
    p.add_argument("inputs", nargs="*", help="Input PPTX files to merge. If omitted, templates are used.")
    p.add_argument("--template", "-t", help="Optional template/master PPTX to use as the base for the merged presentation")
    args = p.parse_args()

    ensure_dir(MERGED_DIR)

    # Resolve input files
    input_files = args.inputs or []
    # If no explicit inputs were provided, gather templates from the templates dir
    if not input_files:
        if os.path.isdir(TEMPLATES_DIR):
            input_files = [os.path.join(TEMPLATES_DIR, n) for n in os.listdir(TEMPLATES_DIR) if n.lower().endswith(".pptx")]

    # If a template was explicitly provided via CLI, set it and ensure it exists
    cli_template = None
    if args.template:
        if os.path.isfile(args.template):
            cli_template = args.template
        else:
            print(f"Template file not found: {args.template}")
            sys.exit(1)

    if not input_files:
        print("No input presentations found. Provide files or add templates to assets/slides/templates.")
        sys.exit(1)

    output_path = os.path.join(MERGED_DIR, args.output)
    # If the user passed a template, use it as the base presentation
    if cli_template:
        # call merge with the template as the first element but signal that
        # merged should be initialized from that template by prepending it
        # to the input_files list but not duplicating when copying.
        input_files = [cli_template] + input_files
    merge_presentations(output_path, input_files)


if __name__ == "__main__":
    main()
