#!/usr/bin/env python3
"""Merge PPTX files by converting every slide to an image and composing
those images into a new PPTX. This guarantees visual fidelity but makes
slides non-editable.

Requires `soffice` (LibreOffice) installed for conversion. The script
will detect `soffice` or `libreoffice` on PATH and use it for slide-to-PNG
conversion.

Usage:
  python3 src/merge_images.py output.pptx input1.pptx input2.pptx ...
"""
import argparse
import os
import shlex
import shutil
import subprocess
import sys
import tempfile
from pptx import Presentation


def find_soffice():
    # Check common names on PATH first
    for cmd in ("soffice", "libreoffice"):
        path = shutil.which(cmd)
        if path:
            return path

    # On macOS, LibreOffice may be installed as an app without adding soffice to PATH.
    # Check the standard application bundle location.
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_path):
        return mac_path

    return None


def convert_to_png(soffice_cmd, src, outdir):
    # Use LibreOffice headless convert-to png
    cmd = [soffice_cmd, "--headless", "--convert-to", "png", "--outdir", outdir, src]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


def build_presentation_from_images(image_paths, output_path):
    prs = Presentation()
    # try to use a blank layout if available
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for img in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Fit the image to the slide size
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(output_path)


def main():
    p = argparse.ArgumentParser(description="Merge PPTX files by rendering slides to images")
    p.add_argument("output", help="Output filename (written to assets/slides/merged)")
    p.add_argument("inputs", nargs="+", help="Input PPTX files to convert and merge")
    args = p.parse_args()

    soffice = find_soffice()
    if not soffice:
        print("Error: LibreOffice (`soffice` or `libreoffice`) not found on PATH. Install LibreOffice to use this option.")
        sys.exit(1)

    root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    merged_dir = os.path.join(root, "assets", "slides", "merged")
    os.makedirs(merged_dir, exist_ok=True)
    output_path = os.path.join(merged_dir, args.output)

    image_list = []
    with tempfile.TemporaryDirectory() as tmp:
        for src in args.inputs:
            if not os.path.isfile(src):
                print(f"Warning: input not found: {src}")
                continue
            convert_to_png(soffice, src, tmp)
            # LibreOffice emits files named after the source base; collect them
            base = os.path.splitext(os.path.basename(src))[0]
            # gather files that start with base and end with .png, sort to keep slide order
            matches = [os.path.join(tmp, f) for f in os.listdir(tmp) if f.startswith(base) and f.lower().endswith('.png')]
            matches.sort()
            image_list.extend(matches)

        if not image_list:
            print("No images produced; aborting.")
            sys.exit(1)

        build_presentation_from_images(image_list, output_path)
        print(f"Saved merged image-based presentation to: {output_path}")


if __name__ == "__main__":
    main()
